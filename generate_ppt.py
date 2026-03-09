"""
Script to generate the Intelli-Doc AI project presentation.
Run: python generate_ppt.py
Output: Intelli-Doc-AI-Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x11, 0x17)   # very dark navy
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)   # sky-blue (primary)
ACCENT2   = RGBColor(0x34, 0xD3, 0x99)   # emerald-green (secondary)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xCB, 0xD5, 0xE1)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)   # card background
MID_TEXT  = RGBColor(0x94, 0xA3, 0xB8)   # muted subtitle text


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Slide helpers ──────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     line_color=None, corner_size=Emu(914400 // 8)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def multiline_textbox(slide, lines, left, top, width, height,
                      font_size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add a textbox with multiple paragraphs (list of strings)."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 – Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # top accent bar
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)

    # left decorative bar
    add_rect(slide, 0, 0.07, 0.08, 7.43, ACCENT)

    # Institution label
    add_textbox(slide, "Atma Malik Institute of Technology and Research (AMRIT)",
                0.2, 0.15, 9.6, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Department of Computer Engineering  ·  University of Mumbai",
                0.2, 0.52, 9.6, 0.35, font_size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)

    # Horizontal divider
    add_rect(slide, 0.6, 0.92, 8.8, 0.025, ACCENT)

    # Main title
    add_textbox(slide, "Intelli-Doc AI",
                0.6, 1.05, 8.8, 1.0, font_size=48, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "An Agentic AI-Powered Code Documentation Platform",
                0.6, 1.95, 8.8, 0.65, font_size=22, bold=False,
                color=ACCENT, align=PP_ALIGN.CENTER)

    # Horizontal divider
    add_rect(slide, 0.6, 2.72, 8.8, 0.025, ACCENT2)

    # Guide label
    add_textbox(slide, "Guide",
                0.6, 2.9, 2.5, 0.4, font_size=12, bold=True,
                color=ACCENT2, align=PP_ALIGN.LEFT)
    add_textbox(slide, "Prof. Swati Bhoir",
                0.6, 3.22, 4.5, 0.45, font_size=16, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(slide, "Dept. of Computer Engineering, AMRIT",
                0.6, 3.6, 5.0, 0.35, font_size=11, color=MID_TEXT, align=PP_ALIGN.LEFT)

    # Team label
    add_textbox(slide, "Team Members",
                5.5, 2.9, 4.0, 0.4, font_size=12, bold=True,
                color=ACCENT, align=PP_ALIGN.LEFT)
    members = ["Rishi Padala", "Sarvesh Pal", "Sarang Patil", "David Kumar"]
    for i, m in enumerate(members):
        add_textbox(slide, f"• {m}",
                    5.5, 3.22 + i * 0.42, 4.0, 0.4, font_size=14, bold=False,
                    color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # Horizontal divider
    add_rect(slide, 0.6, 5.5, 8.8, 0.025, ACCENT)

    # Tags
    tags = ["AI", "Spring Boot", "React", "MongoDB", "Kafka", "Redis", "Gemini 2.5 Flash"]
    x = 0.6
    for t in tags:
        w = len(t) * 0.12 + 0.6
        add_rounded_rect(slide, x, 5.6, w, 0.38, DARK_CARD, line_color=ACCENT)
        add_textbox(slide, t, x + 0.08, 5.65, w - 0.1, 0.32,
                    font_size=11, color=ACCENT, align=PP_ALIGN.CENTER)
        x += w + 0.15

    # bottom accent bar
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)


def slide_introduction(prs):
    """Slide 2 – Introduction / Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    add_textbox(slide, "Introduction", 0.4, 0.18, 9.0, 0.55,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.75, 1.2, 0.05, ACCENT)

    body = (
        "Writing documentation is one of the most time-consuming — yet critical — "
        "parts of software development. Developers often skip it, leading to "
        "poorly maintained codebases that are hard to onboard and understand.\n\n"
        "Intelli-Doc AI solves this by automatically analysing a GitHub repository "
        "and generating structured, human-readable documentation using Google "
        "Gemini 2.5 Flash — in seconds, not hours."
    )
    add_textbox(slide, body, 0.4, 0.95, 9.2, 1.8,
                font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # Three "stat" cards
    cards = [
        ("⚡", "Instant", "Documentation generated\nin seconds"),
        ("🤖", "Agentic AI", "Multi-step AI pipeline\nfor deep understanding"),
        ("🔗", "GitHub Native", "Paste a GitHub URL —\nwe handle the rest"),
    ]
    for i, (icon, title, desc) in enumerate(cards):
        cx = 0.4 + i * 3.15
        add_rounded_rect(slide, cx, 3.0, 2.9, 2.8, DARK_CARD, line_color=ACCENT)
        add_textbox(slide, icon, cx + 1.1, 3.15, 0.8, 0.5, font_size=22, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, cx + 0.1, 3.75, 2.7, 0.5,
                    font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc, cx + 0.1, 4.3, 2.7, 0.9,
                    font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "2", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 3 – Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    add_textbox(slide, "Problem Statement", 0.4, 0.18, 9.0, 0.55,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.75, 1.5, 0.05, ACCENT)

    problems = [
        ("📄", "No Auto-Documentation",
         "Developers write code but rarely document it, making maintenance and "
         "onboarding extremely difficult."),
        ("⏳", "Time-Consuming Process",
         "Manually documenting a mid-sized codebase can take days of developer time."),
        ("🔍", "Inconsistent Quality",
         "Documentation quality varies wildly across contributors and time."),
        ("🔄", "Documentation Drift",
         "Even when docs exist, they quickly become outdated as code evolves."),
    ]
    for i, (icon, title, desc) in enumerate(problems):
        row = i // 2
        col = i % 2
        cx = 0.4 + col * 4.75
        cy = 1.05 + row * 3.0
        add_rounded_rect(slide, cx, cy, 4.45, 2.7, DARK_CARD, line_color=ACCENT)
        add_textbox(slide, icon, cx + 0.2, cy + 0.15, 0.6, 0.5, font_size=20)
        add_textbox(slide, title, cx + 0.85, cy + 0.18, 3.4, 0.45,
                    font_size=15, bold=True, color=ACCENT)
        add_textbox(slide, desc, cx + 0.2, cy + 0.72, 4.0, 1.7,
                    font_size=12, color=LIGHT_GRAY)

    add_textbox(slide, "3", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    """Slide 4 – System Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    add_textbox(slide, "System Architecture", 0.4, 0.18, 9.0, 0.55,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.75, 1.7, 0.05, ACCENT)

    layers = [
        (ACCENT,  "Frontend  (React + TypeScript)",
         "Auth · Dashboard · Repo Detail · Markdown Editor · File Tree · WebSocket"),
        (ACCENT2, "Backend  (Spring Boot 3 · Java 17)",
         "JWT Auth · JGit Clone · Kafka Queue · Redis Cache · WebSocket Notifier"),
        (RGBColor(0xA7, 0x8B, 0xFA), "AI Service  (FastAPI · Python)",
         "Gemini 2.5 Flash · File Selector (Architect) · Doc Generator · Key Rotation"),
        (RGBColor(0xFB, 0xBF, 0x24), "Data Layer  (MongoDB · Redis · Kafka)",
         "User/Repo/Doc persistence · Documentation cache · Async processing queue"),
    ]
    for i, (color, title, desc) in enumerate(layers):
        cy = 1.0 + i * 1.48
        add_rounded_rect(slide, 0.4, cy, 9.2, 1.35, DARK_CARD, line_color=color)
        add_rect(slide, 0.4, cy, 0.12, 1.35, color)
        add_textbox(slide, title, 0.65, cy + 0.12, 8.7, 0.45,
                    font_size=14, bold=True, color=color)
        add_textbox(slide, desc, 0.65, cy + 0.58, 8.7, 0.65,
                    font_size=11, color=LIGHT_GRAY)
        if i < len(layers) - 1:
            add_textbox(slide, "▼", 4.8, cy + 1.3, 0.5, 0.25,
                        font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)

    add_textbox(slide, "4", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


def slide_tech_stack(prs):
    """Slide 5 – Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    add_textbox(slide, "Technology Stack", 0.4, 0.18, 9.0, 0.55,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.75, 1.5, 0.05, ACCENT)

    cols = [
        (ACCENT, "⚛  Frontend",
         ["React 18 + TypeScript", "Vite 5 (bundler)", "Tailwind CSS 3",
          "Radix UI components", "Zustand (state)", "TanStack Query", "React Router v6"]),
        (ACCENT2, "🍃  Backend",
         ["Spring Boot 3.5", "Java 17", "Spring Security + JWT",
          "JGit (repo cloning)", "Retrofit + WebClient", "Springdoc OpenAPI"]),
        (RGBColor(0xA7, 0x8B, 0xFA), "🤖  AI Service",
         ["FastAPI (Python 3.11)", "Google Gemini 2.5 Flash",
          "Multi-key rotation", "Token-bucket rate limiter"]),
        (RGBColor(0xFB, 0xBF, 0x24), "🗄  Data & Infra",
         ["MongoDB (persistence)", "Redis (cache)", "Apache Kafka (queue)",
          "Docker + Compose", "Vercel (frontend)"]),
    ]
    for i, (color, heading, items) in enumerate(cols):
        cx = 0.25 + i * 2.42
        add_rounded_rect(slide, cx, 1.0, 2.2, 6.3, DARK_CARD, line_color=color)
        add_textbox(slide, heading, cx + 0.1, 1.1, 2.0, 0.45,
                    font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_rect(slide, cx + 0.15, 1.57, 1.9, 0.03, color)
        for j, item in enumerate(items):
            add_textbox(slide, f"• {item}", cx + 0.15, 1.72 + j * 0.6, 1.95, 0.55,
                        font_size=11, color=LIGHT_GRAY)

    add_textbox(slide, "5", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


def slide_features(prs):
    """Slide 6 – Key Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    add_textbox(slide, "Key Features", 0.4, 0.18, 9.0, 0.55,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.75, 1.1, 0.05, ACCENT)

    features = [
        ("🔐", "GitHub OAuth + Email Auth", "Secure authentication with OTP email verification"),
        ("📁", "Repo Analysis", "Clones & analyses any public GitHub repository via JGit"),
        ("🧠", "Agentic AI Pipeline", "2-step AI: file selection then deep doc generation"),
        ("⚡", "Real-time Updates", "Kafka + WebSocket live processing progress stream"),
        ("📝", "Markdown Editor", "Edit and customise generated documentation inline"),
        ("🔍", "Full-text Search", "Search across all generated documentation instantly"),
        ("🌲", "File Tree View", "Visual project structure with syntax-highlighted files"),
        ("💾", "Redis Caching", "Instant doc retrieval — no repeated AI calls"),
        ("📊", "Dashboard Stats", "Analytics on repos, docs generated, and usage"),
    ]
    for i, (icon, title, desc) in enumerate(features):
        row = i // 3
        col = i % 3
        cx = 0.35 + col * 3.1
        cy = 1.0 + row * 2.1
        add_rounded_rect(slide, cx, cy, 2.9, 1.9, DARK_CARD, line_color=ACCENT)
        add_textbox(slide, icon, cx + 0.15, cy + 0.1, 0.5, 0.45, font_size=18)
        add_textbox(slide, title, cx + 0.7, cy + 0.12, 2.0, 0.45,
                    font_size=12, bold=True, color=ACCENT)
        add_textbox(slide, desc, cx + 0.15, cy + 0.65, 2.65, 1.1,
                    font_size=10, color=LIGHT_GRAY)

    add_textbox(slide, "6", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


def slide_workflow(prs):
    """Slide 7 – How It Works (Workflow)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    add_textbox(slide, "How It Works", 0.4, 0.18, 9.0, 0.55,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.75, 1.1, 0.05, ACCENT)

    steps = [
        (ACCENT,  "1", "User Authentication",
         "Sign up via email (OTP verification) or one-click GitHub OAuth"),
        (ACCENT2, "2", "Add Repository",
         "Paste any public GitHub URL — backend queues it via Kafka"),
        (ACCENT,  "3", "Clone & Analyse",
         "JGit clones the repo; backend maps the complete file tree"),
        (ACCENT2, "4", "AI File Selection",
         "Gemini AI (Architect) selects the most relevant source files"),
        (ACCENT,  "5", "Generate Docs",
         "Gemini 2.5 Flash generates structured markdown documentation"),
        (ACCENT2, "6", "Cache & Display",
         "Docs cached in Redis; displayed with live search and editor"),
    ]
    for i, (color, num, title, desc) in enumerate(steps):
        row = i // 2
        col = i % 2
        cx = 0.35 + col * 4.8
        cy = 1.05 + row * 2.05
        add_rounded_rect(slide, cx, cy, 4.55, 1.9, DARK_CARD, line_color=color)
        add_rounded_rect(slide, cx + 0.15, cy + 0.55, 0.55, 0.55,
                         color, line_color=None)
        add_textbox(slide, num, cx + 0.15, cy + 0.6, 0.55, 0.45,
                    font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, cx + 0.85, cy + 0.18, 3.5, 0.45,
                    font_size=14, bold=True, color=color)
        add_textbox(slide, desc, cx + 0.85, cy + 0.65, 3.5, 1.0,
                    font_size=11, color=LIGHT_GRAY)

    add_textbox(slide, "7", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


def slide_results(prs):
    """Slide 8 – Results & Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    add_textbox(slide, "Results & Impact", 0.4, 0.18, 9.0, 0.55,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.75, 1.4, 0.05, ACCENT)

    # Metric cards
    metrics = [
        (ACCENT,  "< 60 sec",  "End-to-end documentation\ngeneration time"),
        (ACCENT2, "2-Step AI", "Architect selects files;\nGenerator writes docs"),
        (RGBColor(0xA7, 0x8B, 0xFA), "100 %",  "Automated — zero manual\ndoc effort required"),
    ]
    for i, (color, val, desc) in enumerate(metrics):
        cx = 0.4 + i * 3.1
        add_rounded_rect(slide, cx, 1.0, 2.9, 1.8, DARK_CARD, line_color=color)
        add_textbox(slide, val, cx + 0.1, 1.1, 2.7, 0.7,
                    font_size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc, cx + 0.1, 1.8, 2.7, 0.85,
                    font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Key Achievements", 0.4, 3.0, 9.0, 0.45,
                font_size=16, bold=True, color=WHITE)
    achievements = [
        "✅  Full-stack production deployment on Vercel (frontend) and Docker (backend/AI)",
        "✅  Agentic pipeline using Gemini 2.5 Flash with multi-key rotation & rate limiting",
        "✅  Real-time progress notifications via Kafka + WebSocket integration",
        "✅  Redis caching layer eliminates redundant AI calls for the same repository",
        "✅  Supports public GitHub repos of any language or framework",
        "✅  Seamless developer experience — from URL to full docs in under a minute",
    ]
    for i, a in enumerate(achievements):
        add_textbox(slide, a, 0.4, 3.5 + i * 0.58, 9.2, 0.52,
                    font_size=12, color=LIGHT_GRAY)

    add_textbox(slide, "8", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


def slide_future(prs):
    """Slide 9 – Future Scope"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    add_textbox(slide, "Future Scope", 0.4, 0.18, 9.0, 0.55,
                font_size=30, bold=True, color=WHITE)
    add_rect(slide, 0.4, 0.75, 1.1, 0.05, ACCENT)

    future_items = [
        ("🔒", "Private Repository Support",
         "OAuth token passthrough to access private GitHub repos securely"),
        ("💬", "AI Chat on Codebase",
         "Ask questions about any repository using a RAG-powered chatbot"),
        ("🔔", "Auto-Sync Webhooks",
         "GitHub webhooks to auto-regenerate docs on every push"),
        ("🌐", "Multi-Language Localisation",
         "Generate documentation in multiple human languages"),
        ("🧩", "IDE Plugins",
         "VS Code / JetBrains extension for in-editor documentation"),
        ("📦", "Export to PDF / Confluence",
         "One-click export to PDF, Notion, or Confluence pages"),
    ]
    for i, (icon, title, desc) in enumerate(future_items):
        row = i // 2
        col = i % 2
        cx = 0.35 + col * 4.8
        cy = 1.0 + row * 2.1
        add_rounded_rect(slide, cx, cy, 4.55, 1.9, DARK_CARD,
                         line_color=ACCENT2 if col == 0 else ACCENT)
        add_textbox(slide, icon, cx + 0.2, cy + 0.55, 0.55, 0.5, font_size=22)
        add_textbox(slide, title, cx + 0.85, cy + 0.15, 3.5, 0.45,
                    font_size=13, bold=True,
                    color=ACCENT2 if col == 0 else ACCENT)
        add_textbox(slide, desc, cx + 0.85, cy + 0.65, 3.5, 1.0,
                    font_size=11, color=LIGHT_GRAY)

    add_textbox(slide, "9", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


def slide_thank_you(prs):
    """Slide 10 – Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 0.07, ACCENT)
    add_rect(slide, 0, 7.43, 10, 0.07, ACCENT2)

    # Large decorative circle (background)
    from pptx.util import Emu
    circle = slide.shapes.add_shape(9, Inches(3.0), Inches(1.5), Inches(4.0), Inches(4.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_CARD
    circle.line.fill.background()

    add_textbox(slide, "Thank You!", 0.4, 1.8, 9.2, 1.2,
                font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Intelli-Doc AI  ·  An Agentic AI-Powered Code Documentation Platform",
                0.4, 3.1, 9.2, 0.55, font_size=15, color=ACCENT, align=PP_ALIGN.CENTER)

    add_rect(slide, 2.5, 3.75, 5.0, 0.04, ACCENT2)

    add_textbox(slide, "🌐  https://intelli-docai.vercel.app/",
                1.0, 4.0, 8.0, 0.45, font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Guide  ·  Prof. Swati Bhoir",
                1.0, 4.6, 8.0, 0.45, font_size=13, color=ACCENT2, align=PP_ALIGN.CENTER)

    members_line = "Rishi Padala   ·   Sarvesh Pal   ·   Sarang Patil   ·   David Kumar"
    add_textbox(slide, members_line, 0.4, 5.1, 9.2, 0.45,
                font_size=13, color=MID_TEXT, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Dept. of Computer Engineering  ·  AMRIT  ·  University of Mumbai",
                0.4, 5.65, 9.2, 0.45, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)

    add_textbox(slide, "10", 9.5, 7.1, 0.4, 0.4, font_size=11, color=MID_TEXT, align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_introduction(prs)
    slide_problem(prs)
    slide_architecture(prs)
    slide_tech_stack(prs)
    slide_features(prs)
    slide_workflow(prs)
    slide_results(prs)
    slide_future(prs)
    slide_thank_you(prs)

    out = "Intelli-Doc-AI-Presentation.pptx"
    prs.save(out)
    print(f"✅  Saved → {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
